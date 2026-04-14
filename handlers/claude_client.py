"""
Работа с Claude API: анализ презентаций и генерация улучшений.
Включает полный стилевой промпт AMA Private Club.
"""
import json
import logging
import anthropic
from config import ANTHROPIC_API_KEY

logger = logging.getLogger(__name__)
client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

# ──────────────────────────────────────────────
# СТИЛЕВОЙ ПРОМПТ
# ──────────────────────────────────────────────
STYLE_PROMPT = """
# Промт для создания PPTX-презентаций
## Формат 16:9 · Стиль AMA Private Club

## ФОРМАТ И РАЗМЕР
- Слайды: 16:9, размер 33.87 × 19.05 см (1920 × 1080 px эквивалент)
- Поля: минимум 1.4 см со всех сторон, рабочая зона 31 × 16.2 см
- Один слайд — одна мысль. Не пытаться вместить всё.

## ШРИФТЫ
Заголовки H1/H2: OmegaUIGeometric Bold (700). Если недоступен — Neue Haas Grotesk / Aktiv Grotesk / GT Walsheim.
Подзаголовки: OmegaUIGeometric Medium (500).
Тело: OmegaUI Regular (400). Если недоступен — Inter / Helvetica Neue.
Подписи/лейблы: OmegaUI Medium (500), uppercase, letter-spacing +0.22em.

## ТИПОГРАФИКА — РАЗМЕРЫ
H1 (hero): 54-72pt, межстрочный 0.96-1.0, letter-spacing -0.02em
H2 (slide title): 38-52pt, межстрочный 1.0-1.1, letter-spacing -0.02em
H3 (card title): 18-22pt, межстрочный 1.1, letter-spacing -0.02em
Тело абзаца: 14-17pt, межстрочный 1.75-1.8
Описание (muted): 12-14pt, межстрочный 1.55-1.6
Лейбл (капс): 10-12pt, letter-spacing +0.22em
Футер / копирайт: 9-11pt

## ТИПОГРАФИКА — ПРАВИЛА НАБОРА (РУССКИЙ ТЕКСТ)
- Точки ставятся между предложениями в заголовке как разделитель. НЕ ставятся в конце заголовка и буллета.
- Неразрывные пробелы после предлогов и союзов.
- Нерастяжимый дефис (U+2011) в составных словах.
- Запрещено: предлог/союз в конце строки, точка в конце заголовка/буллета.

## ДОСТУПНЫЕ ЛЕЙАУТЫ — ОБЯЗАТЕЛЬНО ИСПОЛЬЗУЙ РАЗНООБРАЗИЕ

### layout: "hero" (ТОЛЬКО для слайда 1 — обязательно!)
Титульный слайд. Крупный заголовок (54-72pt), подзаголовок, декоративные полосы.
Поле body: подзаголовок/слоган. НЕ добавлять label — они выглядят шаблонно.
ВСЕГДА ПЕРВЫЙ СЛАЙД = hero.

### layout: "quote" или "insight" (слайды с цитатами/инсайтами)
Крупная цитата или ключевой инсайт. Акцент на одной мысли. label сверху.

### layout: "2col" (два столбца — философия, сравнение)
Левая 55%: заголовок + тело. Правая 45%: columns[] со статистикой.
columns: [{"value": "12+", "description": "лет опыта"}, ...]

### layout: "3col" (три столбца — преимущества, этапы)
Три карточки рядом. Каждая: number, title, description.
columns: [{"number": "01", "title": "...", "description": "..."}, ...]

### layout: "4col" (четыре столбца — процесс, шаги)
Четыре шага/блока. Над ними линия с номерами. Каждый: number, title, description.
columns: [{"number": "01", "title": "...", "description": "..."}, ...]

### layout: "6col" (шесть столбцов — услуги, инфраструктура)
Шесть узких карточек. Каждая: title, description. С акцентной линией сверху.
columns: [{"title": "...", "description": "..."}, ...]

### layout: "cta" (ТОЛЬКО для последнего слайда — призыв к действию)
Заголовок по центру + контакты/призыв. ВСЕГДА ПОСЛЕДНИЙ СЛАЙД = cta.

### layout: "content" (универсальный — если не подходит другое)
Лейбл + крупный заголовок + тело текста.

## ПРАВИЛО РАЗНООБРАЗИЯ ЛЕЙАУТОВ
- Никогда не повторять один layout два раза подряд
- Обязательный минимум: hero (1й), quote или insight, минимум один 3col/4col/6col, cta (последний)
- Чередовать: hero → quote → 2col → 4col → content → 3col → cta

## СТРУКТУРА ТИПИЧНОЙ ПРЕЗЕНТАЦИИ (количество слайдов — по контенту)
1. hero — главный слоган (ОБЯЗАТЕЛЕН)
2. quote/insight — инсайт или цитата
3. 2col — философия/подход
4. 4col — процесс/шаги
5. 3col — преимущества
6. 6col — услуги/инфраструктура
7. content — дополнительный блок
8. 2col — почему мы
9. cta — контакты (ОБЯЗаТЕЛЕН)

## ЗАПРЕЩЕНО
- Одинаковый layout на слайдах подряд
- Первый слайд НЕ hero — КРИТИЧЕСКАЯ ОШИБКА
- Последний слайд НЕ cta — КРИТИЧЕСКАЯ ОШИБКА
- Более 5-6 строк текста на слайде
- Bullet points стандартные PowerPoint
- Теги-лейблы на каждом слайде ("INVESTMENT OPPORTUNITY", "KEY INSIGHT" и т.д.) — это маркер AI-генерации
"""

THEME_DARK = """
## ТЕМА — ТЁМНАЯ (Dark)
Все слайды на тёмных фонах, акцент — единственный источник цвета.
Фон основной: #07090E
Фон вторичный: #0C1018
Фон финальный: #040509
Заголовки: #FDFAF5
Тело текста: rgba(241,233,215, 0.65)
Мутный текст: rgba(241,233,215, 0.42)
Акцент: #E08256 (терракота)
Разделители: rgba(241,233,215, 0.06-0.10)
Распределение фонов: 1-Hero #07090E, 2-Инсайт #0C1018, 3-Философия #0C1018, 4-Услуги #07090E, 5-Портфель #0C1018, 6-Инфра #07090E, 7-Процесс #0C1018, 8-Почему #07090E, 9-CTA #040509.
Кнопка primary: фон #E08256, текст #07090E. Лейблы: #E08256. Акцентные цифры: #E08256.
"""

THEME_LIGHT = """
## ТЕМА — СВЕТЛАЯ (Light)
Все слайды на светлом/кремовом фоне. Акцент — тёмно-терракотовый.
Фон основной: #F1E9D7
Фон вторичный: #EAE0CC
Фон финальный: #E0D5C0
Заголовки: #07090E
Тело текста: rgba(10,12,18, 0.65)
Мутный текст: rgba(10,12,18, 0.42)
Акцент: #A05C3A (тёмная терракота)
Разделители: rgba(10,12,18, 0.08-0.12)
Распределение фонов: 1-Hero #F1E9D7, 2-Инсайт #F1E9D7, 3-Философия #EAE0CC, 4-Услуги #F1E9D7, 5-Портфель #EAE0CC, 6-Инфра #F1E9D7, 7-Процесс #EAE0CC, 8-Почему #F1E9D7, 9-CTA #E0D5C0.
Кнопка primary: фон #A05C3A, текст #F1E9D7. Лейблы: #A05C3A. Акцентные цифры: #A05C3A.
"""

THEME_COMBINED = """
## ТЕМА — КОМБИНИРОВАННАЯ (Combined) — по умолчанию
Сэндвич: тёмный открывает, светлые дают передышку, тёмный закрывает.
На тёмных слайдах: фон #07090E/#0C1018/#040509, текст #FDFAF5/rgba(241,233,215,0.65), акцент #E08256.
На светлых слайдах: фон #F1E9D7, текст #07090E/rgba(10,12,18,0.65), акцент #A05C3A.
Распределение фонов (оригинальная схема AMA):
1-Hero #07090E тёмный, 2-Инсайт #F1E9D7 светлый, 3-Философия #0C1018 тёмный, 4-Услуги #07090E тёмный,
5-Портфель #0C1018 тёмный, 6-Инфра #07090E тёмный, 7-Процесс #F1E9D7 светлый, 8-Почему #0C1018 тёмный, 9-CTA #040509 чёрный.
Логика: два светлых слайда (2 и 7) разбивают тёмный массив.
На тёмном слайде — все элементы из тёмной системы. На светлом слайде — все элементы из светлой системы.
"""

THEME_MAP = {
    "dark": THEME_DARK,
    "light": THEME_LIGHT,
    "combined": THEME_COMBINED,
}

# ──────────────────────────────────────────────
# ПУБЛИЧНЫЕ ФУНКЦИИ
# ──────────────────────────────────────────────
async def analyze_and_improve(
    file_path: str | None,
    text_request: str,
    theme: str = "combined",
) -> list[dict]:
    slides_text = ""
    if file_path:
        slides_text = _extract_text_from_pptx(file_path)

    theme_prompt = THEME_MAP.get(theme, THEME_COMBINED)
    system_msg = STYLE_PROMPT + "\n\n" + theme_prompt
    prompt = _build_prompt(slides_text, text_request, theme)

    response = client.messages.create(
        model="claude-opus-4-5",
        max_tokens=6000,
        system=system_msg,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = response.content[0].text
    return _parse_slides_json(raw)


async def generate_iteration_variant(
    slides: list[dict],
    direction: str,
    theme: str = "combined",
) -> list[dict]:
    theme_prompt = THEME_MAP.get(theme, THEME_COMBINED)
    system_msg = STYLE_PROMPT + "\n\n" + theme_prompt
    current_json = json.dumps(slides, ensure_ascii=False, indent=2)

    prompt = (
        f"Вот текущая версия слайдов презентации в JSON:\n\n"
        f"```json\n{current_json}\n```\n\n"
        f"Задача: {direction}\n\n"
        "Сохрани стиль AMA Private Club и все правила типографики.\n"
        "ВАЖНО: сохрани разнообразие layout-ов. Слайд 1 = hero, последний = cta.\n"
        "Верни ТОЛЬКО JSON-массив слайдов в том же формате. "
        "Никакого вступления и объяснений — только JSON."
    )

    response = client.messages.create(
        model="claude-opus-4-5",
        max_tokens=6000,
        system=system_msg,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = response.content[0].text
    return _parse_slides_json(raw)


# ──────────────────────────────────────────────
# ВНУТРЕННИЕ ФУНКЦИИ
# ──────────────────────────────────────────────
def _extract_text_from_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        if texts:
            parts.append(f"=== Слайд {i} ===\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _build_prompt(slides_text: str, text_request: str, theme: str) -> str:
    theme_labels = {"dark": "тёмной", "light": "светлой", "combined": "комбинированной"}
    theme_label = theme_labels.get(theme, "комбинированной")

    base = (
        f"Создай контент презентации в {theme_label} теме стиля AMA Private Club.\n\n"
        "КРИТИЧЕСКИ ВАЖНЫЕ ПРАВИЛА:\n"
        '1. Слайд 1 ВСЕГДА layout=\"hero\" — титульный слайд\n'
        '2. Последний слайд ВСЕГДА layout=\"cta\" — призыв к действию\n'
        "3. Используй РАЗНЫЕ layout для каждого слайда (никогда не повторять подряд)\n"
        "4. Обязательно включи: hero, quote, 2col, 4col или 3col, 6col или 3col, cta\n"
        "5. Для layout 2col/3col/4col/6col — заполняй массив columns[]\n"
        "6. Крупные цифры в columns должны быть короткими: $500M+, EBITDA+, 20M+ (не разбивать!)\n"
        "7. НЕ добавляй поле label — оно не используется\n\n"
        "Принципы контента:\n"
        "- Каждый слайд — одна ключевая мысль\n"
        "- Заголовки — выгоды и результаты, не процессы\n"
        "- Текст лаконичный, убедительный, без воды\n"
        "- Профессиональный тон без канцелярита\n\n"
    )

    if slides_text:
        base += f"Текущий контент презентации:\n\n{slides_text}\n\n"

    if text_request:
        base += f"Дополнительный контекст от автора: {text_request}\n\n"

    base += (
        "Для каждого слайда укажи:\n"
        "- slide_index: номер слайда (начиная с 1)\n"
        "- layout: тип лейаута (hero/quote/insight/2col/3col/4col/6col/cta/content)\n"
        "- slide_type: тип слайда (hero/insight/philosophy/services/portfolio/infrastructure/process/why/cta)\n"
        
        "- title: заголовок\n"
        "- body: основной текст (для 1-колоночных слайдов)\n"
        "- columns: массив колонок (для 2col/3col/4col/6col)\n"
        "- bg_color: hex цвет фона согласно выбранной теме\n\n"
        "ПРИМЕРЫ:\n"
        'Слайд 1: {"slide_index": 1, "layout": "hero", "slide_type": "hero", '
        '"title": "Заголовок", "body": "Подзаголовок", "bg_color": "#07090E"}\n'
        'Слайд 4 (4col): {"slide_index": 4, "layout": "4col", "slide_type": "process", '
        '"title": "Наш процесс", "columns": [{"number": "01", "title": "Анализ", "description": "Текст..."}, '
        '{"number": "02", "title": "Стратегия", "description": "Текст..."}], "bg_color": "#07090E"}\n\n'
        "Верни ТОЛЬКО JSON-массив. Никакого вступления, только JSON.\n"
        "Если нет исходного файла — определи оптимальное количество слайдов исходя из контента (обычно 5-12). Если есть исходный файл — СТРОГО сохрани количество слайдов оригинала."
    )
    return base


def _parse_slides_json(raw: str) -> list[dict]:
    text = raw.strip()
    if "```" in text:
        text = text.split("```")[1]
        if text.startswith("json"):
            text = text[4:]
        text = text.strip()
    start = text.find("[")
    end = text.rfind("]") + 1
    if start == -1 or end == 0:
        logger.error("No JSON array in response: %s", raw[:200])
        return [{"slide_index": 1, "layout": "hero", "title": "Ошибка парсинга", "body": raw[:500]}]
    try:
        data = json.loads(text[start:end])
        # Гарантируем layout на каждом слайде
        for i, slide in enumerate(data):
            if "layout" not in slide:
                stype = slide.get("slide_type", "content")
                slide["layout"] = stype if stype in ("hero", "cta", "quote", "insight", "2col", "3col", "4col", "6col") else "content"
            # Первый слайд всегда hero
            if i == 0:
                slide["layout"] = "hero"
            # Последний слайд всегда cta
            if i == len(data) - 1 and len(data) > 1:
                slide["layout"] = "cta"
        return data
    except json.JSONDecodeError as e:
        logger.error("JSON parse error: %s\nRaw: %s", e, raw[:200])
        return [{"slide_index": 1, "layout": "hero", "title": "Ошибка парсинга JSON", "body": str(e)}]
