"""
Обработка .pptx: создание презентаций с визуальными элементами.
Карточки, акцентные линии, заполнение пространства.
"""
import tempfile
import logging
from pptx import Presentation
from pptx.util import Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)
MARGIN = Cm(1.6)
CONTENT_W = SLIDE_W - 2 * MARGIN
CONTENT_H = SLIDE_H - 2 * MARGIN


def process_for_client(file_path, slides_data, theme="combined"):
    prs = _create_from_scratch(slides_data, theme)
    return _save_tmp(prs)


def _create_from_scratch(slides_data, theme):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for item in slides_data:
        slide = prs.slides.add_slide(blank)
        bg_rgb = _parse_color(item.get("bg_color", "#07090E"))
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_rgb

        bri = (bg_rgb[0]*299 + bg_rgb[1]*587 + bg_rgb[2]*114) / 1000
        dark = bri < 128
        tc = _colors(dark)

        layout = item.get("layout", item.get("slide_type", "content"))

        if layout == "hero":
            _build_hero(slide, item, tc)
        elif layout == "cta":
            _build_cta(slide, item, tc)
        elif layout in ("insight", "quote"):
            _build_quote(slide, item, tc)
        elif layout == "2col":
            _build_2col(slide, item, tc)
        elif layout == "3col":
            _build_3col(slide, item, tc)
        elif layout in ("4col", "process", "why"):
            _build_4col(slide, item, tc)
        elif layout in ("6col", "services", "infrastructure"):
            _build_6col(slide, item, tc)
        else:
            _build_content(slide, item, tc)

    return prs


# ─── LAYOUTS ──────────────────────────────────────────────

def _build_hero(slide, item, tc):
    """Hero: крупный заголовок, подзаголовок, акцентная линия, метрики внизу."""
    # Крупный заголовок — центрирован по вертикали
    title = item.get("title", "")
    tb = slide.shapes.add_textbox(MARGIN, Cm(4), Cm(28), Cm(6))
    _style(tb.text_frame, title, tc["title"], Pt(52), bold=True, ls=-200)

    # Акцентная линия под заголовком
    ln = slide.shapes.add_connector(1, MARGIN, Cm(10.5), MARGIN + Cm(6), Cm(10.5))
    ln.line.color.rgb = tc["accent"]
    ln.line.width = Pt(3)

    # Подзаголовок
    body = item.get("body", "")
    if body:
        tb = slide.shapes.add_textbox(MARGIN, Cm(11.5), Cm(24), Cm(3))
        _style(tb.text_frame, body, tc["muted"], Pt(16))

    # Вертикальная акцентная полоса справа
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(31), Cm(3), Cm(0.4), Cm(13)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = tc["accent"]
    bar.line.fill.background()


def _build_cta(slide, item, tc):
    """CTA: центрированный заголовок, разделитель, кнопка-блок."""
    # Горизонтальный акцентный разделитель сверху
    ln = slide.shapes.add_connector(1, Cm(14), Cm(4), Cm(20), Cm(4))
    ln.line.color.rgb = tc["accent"]
    ln.line.width = Pt(2)

    title = item.get("title", "")
    tb = slide.shapes.add_textbox(Cm(3), Cm(5), Cm(28), Cm(5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    _fmt_run(p, tc["title"], Pt(36), True, -200)

    body = item.get("body", "")
    if body:
        # Кнопка-блок с акцентным фоном
        btn = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Cm(9), Cm(12), Cm(16), Cm(2.5)
        )
        btn.fill.solid()
        btn.fill.fore_color.rgb = tc["accent"]
        btn.line.fill.background()
        tf_btn = btn.text_frame
        tf_btn.word_wrap = True
        tf_btn.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf_btn.paragraphs[0].text = body
        # Текст кнопки — контрастный цвет
        for run in tf_btn.paragraphs[0].runs:
            run.font.color.rgb = tc.get("btn_text", RGBColor(0x07, 0x09, 0x0E))
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.name = "Arial"


def _build_quote(slide, item, tc):
    """Quote: крупный декоративный символ кавычки, цитата, подпись."""
    # Крупная декоративная кавычка
    tb_q = slide.shapes.add_textbox(MARGIN, Cm(2), Cm(6), Cm(5))
    tf_q = tb_q.text_frame
    tf_q.word_wrap = False
    p_q = tf_q.paragraphs[0]
    p_q.text = "\u201c"
    for run in p_q.runs:
        run.font.color.rgb = tc["accent"]
        run.font.size = Pt(120)
        run.font.bold = True
        run.font.name = "Georgia"

    # Цитата / инсайт
    title = item.get("title", "")
    tb = slide.shapes.add_textbox(Cm(3), Cm(6), Cm(26), Cm(6))
    _style(tb.text_frame, title, tc["title"], Pt(28), bold=True, ls=-100)

    # Подпись под цитатой
    body = item.get("body", "")
    if body:
        # Акцентная линия перед подписью
        ln = slide.shapes.add_connector(1, Cm(3), Cm(13), Cm(8), Cm(13))
        ln.line.color.rgb = tc["accent"]
        ln.line.width = Pt(2)
        tb = slide.shapes.add_textbox(Cm(3), Cm(13.8), Cm(24), Cm(2))
        _style(tb.text_frame, body, tc["accent"], Pt(18), bold=True)


def _build_2col(slide, item, tc):
    """2col: левая колонка текст, правая — стат-карточки."""
    # Заголовок
    title = item.get("title", "")
    tb = slide.shapes.add_textbox(MARGIN, Cm(2.5), Cm(16), Cm(4))
    _style(tb.text_frame, title, tc["title"], Pt(30), bold=True, ls=-200)

    # Акцентная линия
    ln = slide.shapes.add_connector(1, MARGIN, Cm(7), MARGIN + Cm(4), Cm(7))
    ln.line.color.rgb = tc["accent"]
    ln.line.width = Pt(2)

    # Тело текста слева
    body = item.get("body", "")
    if body:
        tb = slide.shapes.add_textbox(MARGIN, Cm(8), Cm(15), Cm(8))
        _style(tb.text_frame, body, tc["body"], Pt(14))

    # Правая колонка — стат-карточки
    cols = item.get("columns", [])
    rx = Cm(19)
    rw = Cm(13)
    card_h = Cm(4)
    gap = Cm(0.5)
    y = Cm(2.5)
    for col in cols[:4]:
        val = col.get("value", "")
        desc = col.get("description", col.get("title", ""))
        # Карточка-фон
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, rx, y, rw, card_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = tc["card_bg"]
        card.line.fill.background()
        # Крупная цифра
        tb_v = slide.shapes.add_textbox(rx + Cm(0.8), y + Cm(0.3), Cm(11), Cm(2))
        tf_v = tb_v.text_frame
        tf_v.word_wrap = False
        p_v = tf_v.paragraphs[0]
        p_v.text = str(val)
        _fmt_run(p_v, tc["accent"], Pt(36), True)
        # Описание
        if desc:
            tb_d = slide.shapes.add_textbox(rx + Cm(0.8), y + Cm(2.2), Cm(11), Cm(1.5))
            _style(tb_d.text_frame, desc, tc["muted"], Pt(11))
        y += card_h + gap


def _build_3col(slide, item, tc):
    """3col: заголовок + 3 карточки с фоном на всю нижнюю половину."""
    title = item.get("title", "")
    if title:
        tb = slide.shapes.add_textbox(MARGIN, Cm(2), CONTENT_W, Cm(3))
        _style(tb.text_frame, title, tc["title"], Pt(28), bold=True, ls=-200)

    cols = item.get("columns", [])
    n = min(len(cols), 3)
    if n == 0:
        return
    cw = Cm(9.5)
    gap = Cm(0.7)
    card_top = Cm(6)
    card_h = Cm(11.5)

    for i, col in enumerate(cols[:3]):
        x = MARGIN + i * (cw + gap)
        ct = col.get("title", "")
        cd = col.get("description", col.get("body", ""))
        num = col.get("number", "")

        # Фоновая карточка
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, card_top, cw, card_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = tc["card_bg"]
        card.line.fill.background()

        inner_x = x + Cm(0.8)
        inner_w = cw - Cm(1.6)

        # Номер
        if num:
            tb = slide.shapes.add_textbox(inner_x, card_top + Cm(0.8), inner_w, Cm(1.5))
            _style(tb.text_frame, str(num), tc["accent"], Pt(28), bold=True)
        # Заголовок колонки
        if ct:
            y_ct = card_top + Cm(3) if num else card_top + Cm(0.8)
            tb = slide.shapes.add_textbox(inner_x, y_ct, inner_w, Cm(2.5))
            _style(tb.text_frame, ct, tc["title"], Pt(16), bold=True)
        # Описание
        if cd:
            y_cd = card_top + Cm(5.5) if num else card_top + Cm(3.5)
            tb = slide.shapes.add_textbox(inner_x, y_cd, inner_w, Cm(5))
            _style(tb.text_frame, cd, tc["muted"], Pt(12))


def _build_4col(slide, item, tc):
    """4col: заголовок + 4 карточки с фоном, крупные цифры nowrap."""
    title = item.get("title", "")
    if title:
        tb = slide.shapes.add_textbox(MARGIN, Cm(2), CONTENT_W, Cm(3))
        _style(tb.text_frame, title, tc["title"], Pt(28), bold=True, ls=-200)

    cols = item.get("columns", [])
    n = min(len(cols), 4)
    if n == 0:
        return
    cw = Cm(7.2)
    gap = Cm(0.5)
    card_top = Cm(6)
    card_h = Cm(11.5)

    for i, col in enumerate(cols[:4]):
        x = MARGIN + i * (cw + gap)
        num = col.get("number", col.get("value", ""))
        ct = col.get("title", "")
        cd = col.get("description", col.get("body", ""))

        # Фоновая карточка
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, card_top, cw, card_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = tc["card_bg"]
        card.line.fill.background()

        inner_x = x + Cm(0.6)
        inner_w = cw - Cm(1.2)

        # Номер/цифра — nowrap чтобы $500M+ не разрывалось
        if num:
            tb_n = slide.shapes.add_textbox(inner_x, card_top + Cm(0.8), inner_w, Cm(2.5))
            tf_n = tb_n.text_frame
            tf_n.word_wrap = False
            p_n = tf_n.paragraphs[0]
            p_n.text = str(num)
            # Авто-размер: если короткий — крупный, если длинный — меньше
            sz = Pt(32) if len(str(num)) <= 6 else Pt(22)
            _fmt_run(p_n, tc["accent"], sz, True)
        # Заголовок
        if ct:
            y_ct = card_top + Cm(3.8) if num else card_top + Cm(0.8)
            tb = slide.shapes.add_textbox(inner_x, y_ct, inner_w, Cm(2.5))
            _style(tb.text_frame, ct, tc["title"], Pt(14), bold=True)
        # Описание
        if cd:
            y_cd = card_top + Cm(6.5) if num else card_top + Cm(3.5)
            tb = slide.shapes.add_textbox(inner_x, y_cd, inner_w, Cm(4.5))
            _style(tb.text_frame, cd, tc["muted"], Pt(11))


def _build_6col(slide, item, tc):
    """6col: заголовок + 6 карточек с акцентной линией сверху."""
    title = item.get("title", "")
    if title:
        tb = slide.shapes.add_textbox(MARGIN, Cm(2), CONTENT_W, Cm(3))
        _style(tb.text_frame, title, tc["title"], Pt(28), bold=True, ls=-200)

    cols = item.get("columns", [])
    n = min(len(cols), 6)
    if n == 0:
        return
    cw = Cm(4.7)
    gap = Cm(0.35)
    card_top = Cm(6)
    card_h = Cm(11.5)

    for i, col in enumerate(cols[:6]):
        x = MARGIN + i * (cw + gap)
        ct = col.get("title", "")
        cd = col.get("description", col.get("body", ""))

        # Фоновая карточка
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, card_top, cw, card_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = tc["card_bg"]
        card.line.fill.background()

        # Акцентная линия сверху карточки
        ln = slide.shapes.add_connector(1, x + Cm(0.5), card_top + Cm(0.4), x + Cm(2), card_top + Cm(0.4))
        ln.line.color.rgb = tc["accent"]
        ln.line.width = Pt(2.5)

        inner_x = x + Cm(0.5)
        inner_w = cw - Cm(1)

        if ct:
            tb = slide.shapes.add_textbox(inner_x, card_top + Cm(1.2), inner_w, Cm(2))
            _style(tb.text_frame, ct, tc["title"], Pt(12), bold=True)
        if cd:
            tb = slide.shapes.add_textbox(inner_x, card_top + Cm(3.5), inner_w, Cm(7))
            _style(tb.text_frame, cd, tc["muted"], Pt(10))


def _build_content(slide, item, tc):
    """Content: заголовок, акцентная линия, тело текста."""
    title = item.get("title", "")
    tb = slide.shapes.add_textbox(MARGIN, Cm(3), Cm(28), Cm(4))
    _style(tb.text_frame, title, tc["title"], Pt(30), bold=True, ls=-200)

    # Акцентная линия
    ln = slide.shapes.add_connector(1, MARGIN, Cm(7.5), MARGIN + Cm(4), Cm(7.5))
    ln.line.color.rgb = tc["accent"]
    ln.line.width = Pt(2)

    body = item.get("body", "")
    if body:
        tb = slide.shapes.add_textbox(MARGIN, Cm(8.5), Cm(28), Cm(8))
        _style(tb.text_frame, body, tc["body"], Pt(14))


# ─── HELPERS ──────────────────────────────────────────────

def _colors(dark):
    if dark:
        return {
            "title": RGBColor(0xFD, 0xFA, 0xF5),
            "body": RGBColor(0xBD, 0xB8, 0xAB),
            "muted": RGBColor(0x8A, 0x86, 0x7B),
            "accent": RGBColor(0xE0, 0x82, 0x56),
            "line": RGBColor(0x2A, 0x2D, 0x35),
            "card_bg": RGBColor(0x0F, 0x12, 0x1A),
            "btn_text": RGBColor(0x07, 0x09, 0x0E),
        }
    return {
        "title": RGBColor(0x07, 0x09, 0x0E),
        "body": RGBColor(0x5A, 0x5C, 0x62),
        "muted": RGBColor(0x8A, 0x8C, 0x92),
        "accent": RGBColor(0xA0, 0x5C, 0x3A),
        "line": RGBColor(0xD0, 0xC8, 0xB8),
        "card_bg": RGBColor(0xE2, 0xD9, 0xC5),
        "btn_text": RGBColor(0xF1, 0xE9, 0xD7),
    }


def _parse_color(hex_str):
    h = hex_str.lstrip("#")
    if len(h) != 6:
        return RGBColor(0x07, 0x09, 0x0E)
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except ValueError:
        return RGBColor(0x07, 0x09, 0x0E)


def _style(tf, text, color, size, bold=False, ls=0, spacing=0):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    _fmt_run(p, color, size, bold, ls, spacing)


def _fmt_run(p, color, size, bold=False, ls=0, spacing=0):
    for run in p.runs:
        run.font.color.rgb = color
        run.font.size = size
        run.font.bold = bold
        run.font.name = "Arial"
        if ls:
            run.font._element.attrib['{http://schemas.openxmlformats.org/drawingml/2006/main}spc'] = str(ls)
        if spacing:
            run.font._element.attrib['{http://schemas.openxmlformats.org/drawingml/2006/main}spc'] = str(spacing)


def _save_tmp(prs):
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", prefix="prez_", delete=False)
    prs.save(tmp.name)
    tmp.close()
    return tmp.name
